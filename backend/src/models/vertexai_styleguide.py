"""VertexAI API to create a powerpoint style guide.

Generative AI based on VertexAI's PaLM to create a powerpoint style
guide based on the input information.
"""

import os
import sys
import logging
import json
import re

from langchain.embeddings import VertexAIEmbeddings
from langchain.chains import LLMChain
from langchain.chat_models import ChatVertexAI
from langchain import PromptTemplate
from langchain.output_parsers import StructuredOutputParser
from langchain.output_parsers import ResponseSchema

import collections
import collections.abc
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor


# Path must be defined (e.g. PYTHONPATH="/path/to/repo/backend")
sys.path.append(os.path.abspath("../../../"))

from src.models import ChatBot


class VertexAIStyleGuide(ChatBot):
    """PaLM agent for generating powerpoint style guides."""

    customer_insights = {
        "gender": "",
        "age": "",
        "income": "",
        "style": "",
        "outerwear1": "",
        "outerwear2": "",
        "outerwear3": "",
        "top1": "",
        "top2": "",
        "top3": "",
        "bottom1": "",
        "bottom2": "",
        "bottom3": "",
    }

    def use_default_template(self):
        """Define the default prompt scenario.

        Args:
            None
        Returns:
            None
        """

        # Prepare prompt
        if not self.config.get("prompt", {}).get("identity", False):
            self.identity = "You are a fashion expert working a clothing company."
        else:
            self.identity = self.config["prompt"]["identity"]
        if not self.config.get("prompt", {}).get("intent", False):
            self.intent = "You write marketing material for your clothing company personalized to the customer."
        else:
            self.intent = self.config["prompt"]["intent"]
        if not self.config.get("prompt", {}).get("behavior", False):
            self.behavior = "You are well-spoken and clear."
        else:
            self.behavior = self.config["prompt"]["behavior"]

        # Set defaults
        self.template = self._template()

        return None

    def add_user_input(self, input_dict):
        """Append a text input to the chain.

        Args:
            user_input_dict (dict): Text input to append to chain
        Returns:
            output (dict): Formatted JSON response
        """

        try:
            user_input = f"""
The customer is {input_dict["age"]} years old and {input_dict["gender"]}.
The customer is in the {input_dict["income"]} income bracket group.
The customer prefers the {input_dict["style"]} fashion style.
The customer is considering buying {input_dict["outerwear1"]}, {input_dict["outerwear2"]}, or {input_dict["outerwear3"]} as outerwear.
The customer is considering buying {input_dict["top1"]}, {input_dict["top2"]}, or {input_dict["top3"]} as a top.
The customer is considering buying {input_dict["bottom1"]}, {input_dict["bottom2"]}, or {input_dict["bottom3"]} as bottoms.
"""
            response = self.chain(user_input)

            text = str(response["text"])
            text = re.sub("^```json", "", text)
            text = re.sub("```$", "", text)

            if text:
                output = json.loads(text)
            else:
                output = {}
        except:
            output = {}

        return output

    def _llm(self):
        """Private method to define LangChain LLM.

        Args:
            None
        Returns:
            llm (obj): LangChain LLM object
        """

        api_key = self.config.get("vertexai", {}).get("api_key", "")
        temperature = self.config.get("vertexai", {}).get("temperature", 0.9)
        model_name = self.config.get("vertexai", {}).get("model_name", "chat-bison")
        max_tokens = self.config.get("vertexai", {}).get("max_tokens", "2048")

        llm = ChatVertexAI(
            temperature=temperature,
            model_name=model_name,
            max_output_tokens=max_tokens,
        )

        return llm

    def _embeddings(self):
        """Private method to define LangChain embeddings.

        Args:
            None
        Returns:
            model (obj): Pre-trained model object
        """

        api_key = self.config.get("vertexai", {}).get("api_key", "")
        embeddings = VertexAIEmbeddings()
        return embeddings

    def _chain(self):
        """Private method to define LangChain chain.

        Args:
            None
        Returns:
            chain (obj): LangChain chain object
        """

        chain = LLMChain(
            llm=self.llm,
            prompt=self.prompt,
            verbose=False,
        )
        return chain

    def _memory(self):
        """Private method to define LangChain memory.

        Args:
            None
        Returns:
            memory (obj): LangChain memory object
        """

        memory = None
        return memory

    def _template(self):
        """Private method to define text template.

        Args:
            None
        Returns:
            template (str): Text template for prompt
        """

        template = f"{self.identity}\n{self.intent}\n{self.behavior}\n----------------\n{{format_instructions}}{{question}}\n"
        return template

    def _prompt(self):
        """Private method to define LangChain prompt template.

        Args:
            None
        Returns:
            prompt (obj): LangChain prompt template object
        """

        info = self.customer_insights

        tagline = f"Write a tagline to sell clothing to a {info['age']} year old {info['gender']} customer with {info['income']} income and a preference for {info['style']} fashion."
        brand = f"Make up a brand description for a clothing company tailored for {info['age']} year old {info['gender']} customers with {info['income']} income and a preference for {info['style']} fashion."
        style = f"Write a three sentence style guide containing fashion advice for a customer interested in buying {info['outerwear1']}, {info['top1']}, and {info['bottom1']}."
        collection1 = f"Write a brief two sentence product description in for a collection of containing {info['outerwear1']}, {info['top1']}, and {info['bottom1']}."
        collection2 = f"Write a brief two sentence product description in for a collection of containing {info['outerwear2']}, {info['top2']}, and {info['bottom2']}."
        collection3 = f"Write a brief two sentence product description in for a collection of containing {info['outerwear3']}, {info['top3']}, and {info['bottom3']}."
        testimony1 = f"Make up a one sentence customer testimony for a clothing company from a enthusastic {info['gender']} {info['age']} year old customer who is very enthusiastic and excitable."
        testimony3 = f"Create a different customer testimony for a clothing company from a boring {info['gender']} customer in a {info['income']} income bracket and talks in a mellow way."
        testimony2 = f"Come up with another new customer testimony for a clothing company from a concise {info['gender']} customer who prefers the {info['style']} fashion style, but is very concise and brief."
        conclusion = f"Write one sentence from the perspective of the company personally thanking this specific customer for their time."

        response_schemas = [
            ResponseSchema(
                name="tagline",
                description=tagline,
            ),
            ResponseSchema(
                name="brand",
                description=brand,
            ),
            ResponseSchema(
                name="style",
                description=style,
            ),
            ResponseSchema(
                name="collection1",
                description=collection1,
            ),
            ResponseSchema(
                name="collection2",
                description=collection2,
            ),
            ResponseSchema(
                name="collection3",
                description=collection3,
            ),
            ResponseSchema(
                name="testimony1",
                description=testimony1,
            ),
            ResponseSchema(
                name="testimony2",
                description=testimony2,
            ),
            ResponseSchema(
                name="testimony2",
                description=testimony2,
            ),
            ResponseSchema(
                name="conclusion",
                description=conclusion,
            ),
        ]
        output_parser = StructuredOutputParser.from_response_schemas(response_schemas)

        format_instructions = output_parser.get_format_instructions()
        prompt = PromptTemplate(
            template=self.template,
            input_variables=["question"],
            partial_variables={"format_instructions": format_instructions},
        )

        return prompt

    def generate_powerpoint(self, response, pptx_filename):
        """Generate a powerpoint slide deck based on response.

        Args:
            response (dict): Formatted JSON response
        Returns:
            None
        """

        # Logic for picking a template
        template = "casual"
        if self.customer_insights["style"] in ["formal", "elegant", "business"]:
            template = "formal"

        # Load a template
        if template == "casual":
            slides = Presentation("./data/templates/casual_style_guide.pptx")

            # Title
            slide = slides.slides[0]
            slide.shapes[2].text = response["tagline"]

            # Brand Slide
            slide = slides.slides[2]
            slide.shapes[1].text = response["brand"]

            # Style slide
            slide = slides.slides[3]
            slide.shapes[5].text = response["style"]

            # Collection slide
            slide = slides.slides[4]
            slide.shapes[4].text = response["collection1"]
            slide.shapes[4].text_frame.paragraphs[0].font.size = Pt(12)
            slide.shapes[7].text = response["collection2"]
            slide.shapes[7].text_frame.paragraphs[0].font.size = Pt(12)
            slide.shapes[10].text = response["collection3"]
            slide.shapes[10].text_frame.paragraphs[0].font.size = Pt(12)

            # Testimony slide
            slide = slides.slides[5]
            slide.shapes[2].text = response["testimony1"]
            slide.shapes[2].text_frame.paragraphs[0].font.size = Pt(12)
            slide.shapes[3].text = response["testimony2"]
            slide.shapes[3].text_frame.paragraphs[0].font.size = Pt(12)
            slide.shapes[4].text = response["testimony3"]
            slide.shapes[4].text_frame.paragraphs[0].font.size = Pt(12)

            # Conclusion slide
            slide = slides.slides[6]
            slide.shapes[4].text = response["conclusion"]
            slide.shapes[4].text_frame.paragraphs[0].font.size = Pt(16)

        else:
            slides = Presentation("./data/templates/formal_style_guide.pptx")

            # Title
            slide = slides.slides[0]
            slide.shapes[2].text = response["tagline"]

            # Brand Slide
            slide = slides.slides[2]
            slide.shapes[1].text = response["brand"]

            # Style slide
            slide = slides.slides[3]
            slide.shapes[2].text = response["style"]
            slide.shapes[2].text_frame.paragraphs[0].font.color.rgb = RGBColor(
                0x00, 0x00, 0x00
            )

            # Collection slide
            slide = slides.slides[4]
            slide.shapes[1].text = response["collection1"]
            slide.shapes[1].text_frame.paragraphs[0].font.size = Pt(12)
            slide.shapes[4].text = response["collection2"]
            slide.shapes[4].text_frame.paragraphs[0].font.size = Pt(12)
            slide.shapes[7].text = response["collection3"]
            slide.shapes[7].text_frame.paragraphs[0].font.size = Pt(12)

            # Testimony slide
            slide = slides.slides[5]
            slide.shapes[2].text = response["testimony1"]
            slide.shapes[2].text_frame.paragraphs[0].font.size = Pt(12)
            slide.shapes[3].text = response["testimony2"]
            slide.shapes[3].text_frame.paragraphs[0].font.size = Pt(12)
            slide.shapes[4].text = response["testimony3"]
            slide.shapes[4].text_frame.paragraphs[0].font.size = Pt(12)

            # Conclusion slide
            slide = slides.slides[6]
            slide.shapes[2].text = response["conclusion"]
            slide.shapes[2].text_frame.paragraphs[0].font.size = Pt(16)

        slides.save(pptx_filename)
        return None
